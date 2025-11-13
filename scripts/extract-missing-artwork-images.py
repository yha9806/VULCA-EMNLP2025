#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
从PPT中提取缺失的作品图片
专门为artwork-10和artwork-27提取图片
"""

from pptx import Presentation
from PIL import Image
import io
import os

def extract_images_for_artwork(ppt_path, slides, output_dir, artwork_id):
    """
    从指定的PPT幻灯片中提取图片

    Args:
        ppt_path: PPT文件路径
        slides: 幻灯片编号列表 (从1开始)
        output_dir: 输出目录
        artwork_id: 作品ID (如 'artwork-10')
    """

    print(f"\n正在读取 PPT: {ppt_path}")
    prs = Presentation(ppt_path)

    print(f"共找到 {len(prs.slides)} 张幻灯片")
    print(f"需要提取的幻灯片: {slides}")

    # 创建作品目录
    artwork_dir = os.path.join(output_dir, artwork_id)
    os.makedirs(artwork_dir, exist_ok=True)

    print(f"输出目录: {artwork_dir}\n")

    extracted_count = 0
    image_index = 1

    for slide_num in slides:
        if slide_num > len(prs.slides):
            print(f"  [WARNING] Slide {slide_num} does not exist")
            continue

        slide = prs.slides[slide_num - 1]  # PPT索引从0开始
        print(f"--- Processing Slide {slide_num} ---")

        # 提取这张幻灯片中的所有图片
        slide_image_count = 0
        for shape in slide.shapes:
            if hasattr(shape, "image"):
                try:
                    # 获取图片数据
                    image = shape.image
                    image_bytes = image.blob

                    # 确定图片格式
                    ext = image.ext
                    if ext.startswith('.'):
                        ext = ext[1:]

                    # 保存图片
                    filename = f"{image_index:02d}.{ext}"
                    filepath = os.path.join(artwork_dir, filename)

                    with open(filepath, 'wb') as f:
                        f.write(image_bytes)

                    # 获取图片尺寸
                    img = Image.open(io.BytesIO(image_bytes))
                    width, height = img.size

                    print(f"  [OK] Image {image_index}: {filename} ({width}x{height}px, {len(image_bytes)/1024:.1f}KB)")

                    image_index += 1
                    slide_image_count += 1
                    extracted_count += 1

                except Exception as e:
                    print(f"  [ERROR] Failed to extract image: {e}")

        print(f"  Total images from this slide: {slide_image_count}\n")

    print(f"\n[SUCCESS] Extraction complete!")
    print(f"   Total images extracted: {extracted_count}")
    print(f"   Output directory: {artwork_dir}\n")

    return extracted_count

def main():
    # PPT文件路径
    ppt_path = r"I:\VULCA-EMNLP2025\丛生--沉思之胃 人员作品.pptx"
    output_base = r"I:\VULCA-EMNLP2025\assets\artworks"

    print("="*60)
    print("从PPT中提取缺失的作品图片")
    print("="*60)

    # 提取artwork-10的图片 (PPT 13-14页)
    print("\n[1] 提取 artwork-10 (陈筱薇《源游》)")
    print("    PPT页码: 13-14")
    count1 = extract_images_for_artwork(
        ppt_path=ppt_path,
        slides=[13, 14],
        output_dir=output_base,
        artwork_id='artwork-10'
    )

    # 提取artwork-27的图片 (PPT 51-52页)
    print("\n[2] 提取 artwork-27 (吴振华《辩白书》)")
    print("    PPT页码: 51-52")
    count2 = extract_images_for_artwork(
        ppt_path=ppt_path,
        slides=[51, 52],
        output_dir=output_base,
        artwork_id='artwork-27'
    )

    print("="*60)
    print(f"总结:")
    print(f"  artwork-10: {count1} 张图片")
    print(f"  artwork-27: {count2} 张图片")
    print(f"  总计: {count1 + count2} 张图片")
    print("="*60)

if __name__ == "__main__":
    main()
