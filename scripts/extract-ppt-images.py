#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
从PPT中提取图片文件
将图片保存到展览的assets目录，按作品组织
"""

from pptx import Presentation
from PIL import Image
import io
import os
import json

def extract_images_from_ppt(ppt_path, output_dir, structured_data_path):
    """
    从PPT中提取图片并保存到对应的作品目录

    Args:
        ppt_path: PPT文件路径
        output_dir: 输出目录 (exhibitions/congsheng-2025/assets/artworks/)
        structured_data_path: 结构化作品数据JSON路径
    """

    # 读取结构化作品数据
    with open(structured_data_path, 'r', encoding='utf-8') as f:
        data = json.load(f)

    # 创建幻灯片编号到作品的映射
    slide_to_artwork = {}
    for institution, artworks in data['institutions'].items():
        for artwork in artworks:
            slide_num = artwork['slide_number']
            slide_to_artwork[slide_num] = {
                'artists': artwork['artists'],
                'title': artwork['title'],
                'institution': institution,
                'image_count': artwork['image_count']
            }

    print(f"正在读取 PPT: {ppt_path}")
    prs = Presentation(ppt_path)

    print(f"共找到 {len(prs.slides)} 张幻灯片")
    print(f"共映射 {len(slide_to_artwork)} 件作品\n")

    # 确保输出目录存在
    os.makedirs(output_dir, exist_ok=True)

    extracted_count = 0

    for slide_num, slide in enumerate(prs.slides, start=1):
        # 检查这张幻灯片是否对应一个作品
        if slide_num not in slide_to_artwork:
            continue

        artwork_info = slide_to_artwork[slide_num]

        # 创建作品目录名（使用艺术家名字）
        artists_str = '-'.join(artwork_info['artists'])
        # 清理文件名（移除特殊字符）
        safe_artists = ''.join(c for c in artists_str if c.isalnum() or c in '-_')
        artwork_dir = os.path.join(output_dir, f"artwork-{slide_num:02d}-{safe_artists}")
        os.makedirs(artwork_dir, exist_ok=True)

        print(f"--- Slide {slide_num}: {artwork_info['title']} ---")
        print(f"  艺术家: {', '.join(artwork_info['artists'])}")
        print(f"  学院: {artwork_info['institution']}")

        # 提取这张幻灯片中的所有图片
        image_index = 1
        for shape in slide.shapes:
            if hasattr(shape, "image"):
                try:
                    # 获取图片数据
                    image = shape.image
                    image_bytes = image.blob

                    # 确定图片格式
                    ext = image.ext
                    if ext.startswith('.'):
                        ext = ext[1:]

                    # 保存图片
                    filename = f"{image_index:02d}.{ext}"
                    filepath = os.path.join(artwork_dir, filename)

                    with open(filepath, 'wb') as f:
                        f.write(image_bytes)

                    # 获取图片尺寸
                    img = Image.open(io.BytesIO(image_bytes))
                    width, height = img.size

                    print(f"    [OK] Image {image_index}: {filename} ({width}x{height}px, {len(image_bytes)/1024:.1f}KB)")

                    image_index += 1
                    extracted_count += 1

                except Exception as e:
                    print(f"    [ERROR] Failed to extract image: {e}")

        print()

    print(f"\n[OK] Extraction complete! Saved {extracted_count} images")
    print(f"[OK] Output directory: {output_dir}")

def main():
    ppt_path = r"I:\VULCA-EMNLP2025\丛生--沉思之胃 人员作品.pptx"
    output_dir = r"I:\VULCA-EMNLP2025\exhibitions\congsheng-2025\assets\artworks"
    structured_data_path = r"I:\VULCA-EMNLP2025\scripts\exhibition-artworks-structured.json"

    extract_images_from_ppt(ppt_path, output_dir, structured_data_path)

if __name__ == "__main__":
    main()
