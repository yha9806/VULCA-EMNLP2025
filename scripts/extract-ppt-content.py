#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
提取 PPT 内容脚本
从"丛生--沉思之胃 人员作品.pptx"提取展览信息和作品数据
"""

from pptx import Presentation
import json
import sys

def extract_text_from_shape(shape):
    """从形状中提取文本"""
    if hasattr(shape, "text"):
        return shape.text.strip()
    return ""

def extract_slide_content(slide):
    """提取单个幻灯片的内容"""
    content = {
        "texts": [],
        "images": [],
        "tables": []
    }

    for shape in slide.shapes:
        # 提取文本
        text = extract_text_from_shape(shape)
        if text:
            content["texts"].append(text)

        # 检查是否有图片
        if hasattr(shape, "image"):
            try:
                image_info = {
                    "width": shape.width,
                    "height": shape.height,
                    "left": shape.left,
                    "top": shape.top
                }
                content["images"].append(image_info)
            except:
                pass

        # 检查是否有表格
        if shape.has_table:
            table_data = []
            for row in shape.table.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                table_data.append(row_data)
            content["tables"].append(table_data)

    return content

def main():
    ppt_path = r"I:\VULCA-EMNLP2025\丛生--沉思之胃 人员作品.pptx"

    try:
        print(f"正在读取 PPT 文件: {ppt_path}")
        prs = Presentation(ppt_path)

        print(f"\n[OK] PPT read successfully!")
        print(f"  Total slides: {len(prs.slides)}")
        print(f"  Slide size: {prs.slide_width} x {prs.slide_height}\n")

        # 提取所有幻灯片内容
        all_slides = []

        for idx, slide in enumerate(prs.slides, 1):
            print(f"--- 幻灯片 {idx} ---")
            content = extract_slide_content(slide)

            # 打印文本内容
            if content["texts"]:
                for text in content["texts"]:
                    print(f"  {text}")

            # 打印图片信息
            if content["images"]:
                print(f"  [图片数量: {len(content['images'])}]")

            # 打印表格信息
            if content["tables"]:
                print(f"  [表格数量: {len(content['tables'])}]")
                for table in content["tables"]:
                    for row in table:
                        print(f"    {' | '.join(row)}")

            print()

            all_slides.append({
                "slide_number": idx,
                "content": content
            })

        # 保存为 JSON
        output_path = r"I:\VULCA-EMNLP2025\scripts\ppt-extracted-content.json"
        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(all_slides, f, ensure_ascii=False, indent=2)

        print(f"\n[OK] Content saved to: {output_path}")

    except FileNotFoundError:
        print(f"[ERROR] File not found: {ppt_path}")
        sys.exit(1)
    except Exception as e:
        print(f"[ERROR] {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)

if __name__ == "__main__":
    main()
